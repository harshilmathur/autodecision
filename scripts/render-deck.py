#!/usr/bin/env python3
"""
render-deck.py — render a Decision Brief into a McKinsey-styled PPTX.

First iteration: hardcoded slide spec for the Adobe Photoshop UI deprecation
brief at examples/adobe-deprecate-photoshop-ui.md.

Goal of this iteration: nail the visual fidelity (typography, layout grid,
exhibit style) on a fixed example before generalizing to a brief parser +
LLM action-title synthesis pass.

Usage: python3 scripts/render-deck.py [out_path]
       (defaults to /tmp/adobe-deck.pptx)
"""

from __future__ import annotations
import io
import sys
from dataclasses import dataclass
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ---------------------------------------------------------------------------
# McKinsey style constants
# ---------------------------------------------------------------------------

# Colors — eyeballed from the Chile H2 reference deck
NAVY        = RGBColor(0x05, 0x1C, 0x2C)  # chapter dividers, TOC
INK         = RGBColor(0x1A, 0x1A, 0x1A)  # body text
GREY_RULE   = RGBColor(0xC8, 0xCC, 0xD0)  # horizontal rules
GREY_FOOT   = RGBColor(0x6E, 0x73, 0x78)  # footer text
GREY_LIGHT  = RGBColor(0xE9, 0xEC, 0xEF)  # alt row fill
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
MCK_BLUE    = RGBColor(0x05, 0x4A, 0xD8)  # accent / primary
MCK_BLUE_BR = RGBColor(0x21, 0x96, 0xF3)  # bright accent
MCK_BLUE_PL = RGBColor(0x9D, 0xC3, 0xE6)  # pale fill
ACCENT_CORAL= RGBColor(0xE8, 0x9B, 0xC7)  # contrast accent

# Fonts — Georgia is the closest free analog to McKinsey's Bower/Cheltenham
TITLE_FONT  = "Georgia"
BODY_FONT   = "Helvetica"   # falls back to Arial on non-Mac

# Slide geometry — 16:9, 13.333 × 7.5 inches
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(0.55)
MARGIN_R = Inches(0.55)
MARGIN_T = Inches(0.42)
TITLE_BOTTOM = Inches(1.65)   # below action title + rule (more breathing room)
FOOTER_Y = Inches(7.18)
BULLET = "▪"   # McKinsey-style square bullet


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _set_text(frame, text, *, font=BODY_FONT, size=12, bold=False, color=INK,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return p


def _add_textbox(slide, x, y, w, h, text, **kwargs):
    tb = slide.shapes.add_textbox(x, y, w, h)
    _set_text(tb.text_frame, text, **kwargs)
    return tb


def _add_rule(slide, x, y, w, color=GREY_RULE, weight=0.75):
    line = slide.shapes.add_connector(1, x, y, x + w, y)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def _add_filled_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_circle_badge(slide, x, y, size, label, *, fill=NAVY, txt=WHITE):
    """Filled circle with centered text — used for numbered grids ❶ ❷ ❸."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = BODY_FONT
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = txt
    return shape


def _add_status_pill(slide, x, y, w, label, *, color):
    """Subtle status indicator: small colored disc + label.

    Auto-shrinks font for long text so the pill stays one row visually
    even when callers pass status copy beyond ~22 chars.
    """
    disc = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.05),
                                  Inches(0.13), Inches(0.13))
    disc.fill.solid()
    disc.fill.fore_color.rgb = color
    disc.line.fill.background()
    disc.shadow.inherit = False
    # Adaptive font size: 10pt for short status, 9pt for medium, 8.5pt for long
    n = len(label)
    if n <= 22:
        size = 10
    elif n <= 32:
        size = 9
    else:
        size = 8.5
    _add_textbox(slide, x + Inches(0.22), y, w - Inches(0.22), Inches(0.5),
                 label, font=BODY_FONT, size=size, color=color, bold=True)


def _set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_footer(slide, page_num, source=None, brand="Autodecision"):
    if source:
        _add_textbox(slide, MARGIN_L, FOOTER_Y, Inches(9), Inches(0.25),
                     source, font=BODY_FONT, size=8, color=GREY_FOOT)
    # Right footer: small brand mark + brand + page number
    slide.shapes.add_picture(_make_brand_mark_small(),
                             Inches(11.05), Inches(FOOTER_Y.inches - 0.04),
                             width=Inches(0.18), height=Inches(0.18))
    _add_textbox(slide, Inches(10.0), FOOTER_Y, Inches(3.0), Inches(0.25),
                 f"{brand}     {page_num}",
                 font=BODY_FONT, size=8, color=GREY_FOOT, align=PP_ALIGN.RIGHT)


def _add_action_title(slide, title, *, color=INK, size=22, prefix=None):
    """Bold serif action title + thin rule beneath. Optional 'prefix' renders
    as a small uppercase kicker line ABOVE the title."""
    if prefix:
        _add_textbox(slide, MARGIN_L, MARGIN_T - Inches(0.02),
                     Inches(8), Inches(0.25),
                     prefix, font=BODY_FONT, size=9.5, bold=True,
                     color=MCK_BLUE)
        title_y = MARGIN_T + Inches(0.30)
    else:
        title_y = MARGIN_T
    tb = slide.shapes.add_textbox(MARGIN_L, title_y,
                                  SLIDE_W - MARGIN_L - MARGIN_R, Inches(1.05))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = TITLE_FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    p.line_spacing = 1.10
    # Generous rule offset below title
    rule_y = Inches(1.60) + (Inches(0.30) if prefix else Inches(0))
    _add_rule(slide, MARGIN_L, rule_y,
              SLIDE_W - MARGIN_L - MARGIN_R, color=GREY_RULE, weight=0.75)
    return rule_y


def _add_paragraphs(frame, items, *, font=BODY_FONT, size=11, color=INK,
                    bullet=True, leading=1.30, bullet_color=None):
    """Add multiple paragraphs to an existing text frame.

    If bullet_color is given, the bullet character is rendered in that color
    while body text is in `color`."""
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    bc = bullet_color if bullet_color is not None else MCK_BLUE
    for i, it in enumerate(items):
        if i == 0:
            p = frame.paragraphs[0]
        else:
            p = frame.add_paragraph()
        p.line_spacing = leading
        p.space_after = Pt(5)
        if bullet:
            p.text = ""
            br = p.add_run()
            br.text = BULLET + "  "
            br.font.name = font
            br.font.size = Pt(size)
            br.font.color.rgb = bc
            run = p.add_run()
            run.text = it
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = it
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color


# ---------------------------------------------------------------------------
# Slide types
# ---------------------------------------------------------------------------

def add_title_cover(prs, title, subtitle, footer_org, date):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(s, WHITE)

    # Top-left brand mark — small graphic mark + wordmark
    s.shapes.add_picture(_make_brand_mark_small(),
                         MARGIN_L, Inches(0.48),
                         width=Inches(0.36), height=Inches(0.36))
    _add_textbox(s, Inches(MARGIN_L.inches + 0.45), Inches(0.45),
                 Inches(4), Inches(0.7),
                 footer_org, font=TITLE_FONT, size=20, bold=True, color=NAVY)

    # Cover graphic — the autodecision "Council to Synthesis" mark.
    # Five colored persona nodes in a pentagon converge via gradient
    # lines on a central synthesis node. The mark IS the methodology.
    s.shapes.add_picture(_make_council_mark(),
                         Inches(6.5), Inches(0.7),
                         width=Inches(6.4), height=Inches(6.4))

    # Title — left half of the cover only; right half is the council mark.
    # Width capped at 5.8" so long titles wrap before reaching the graphic
    # (which starts at 6.5"), keeping a clean gutter on every cover. Font
    # adapts to title length so long titles don't break into >3 lines.
    title_tb = s.shapes.add_textbox(MARGIN_L, Inches(2.4), Inches(5.8), Inches(3.4))
    tf = title_tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = TITLE_FONT
    # Pick a font that lets each line of the title fit in the 5.8" box.
    # Width budget per char (Georgia bold): ~0.18" at 38pt, ~0.15" at 32pt,
    # ~0.13" at 28pt. Decision made on the longest line (split on \n) so
    # explicit breaks don't get punished for total length.
    longest_line = max((len(line) for line in title.split("\n")), default=0)
    if longest_line <= 18:
        run.font.size = Pt(38)
    elif longest_line <= 30:
        run.font.size = Pt(32)
    else:
        run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = INK
    p.line_spacing = 1.05

    # Subtitle and date pinned to bottom band, won't collide with title
    _add_textbox(s, MARGIN_L, Inches(6.05), Inches(7), Inches(0.4),
                 subtitle, font=BODY_FONT, size=14, color=INK)
    _add_textbox(s, MARGIN_L, Inches(6.45), Inches(7), Inches(0.35),
                 date, font=BODY_FONT, size=12, color=GREY_FOOT)

    # Confidentiality / provenance footer
    _add_textbox(s, MARGIN_L, Inches(7.05), Inches(10), Inches(0.4),
                 "GENERATED BY AUTODECISION — "
                 "possibility map with synthesized recommendation, not a single-shot answer",
                 font=BODY_FONT, size=8, color=GREY_FOOT)


def add_toc_dark(prs, items, page_num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, NAVY)
    _add_textbox(s, MARGIN_L, Inches(0.55), Inches(8), Inches(0.6),
                 "Table of contents", font=TITLE_FONT, size=28, bold=True, color=WHITE)
    _add_rule(s, MARGIN_L, Inches(1.20), SLIDE_W - MARGIN_L - MARGIN_R,
              color=RGBColor(0x4A, 0x66, 0x7E), weight=0.5)
    # Header row
    _add_textbox(s, MARGIN_L, Inches(1.35), Inches(8), Inches(0.4),
                 "Section", font=BODY_FONT, size=13, color=WHITE)
    _add_textbox(s, Inches(10.5), Inches(1.35), Inches(2.5), Inches(0.4),
                 "Slides", font=BODY_FONT, size=13, color=WHITE)
    _add_rule(s, MARGIN_L, Inches(1.75), SLIDE_W - MARGIN_L - MARGIN_R,
              color=RGBColor(0x4A, 0x66, 0x7E), weight=0.5)

    y = Inches(1.95)
    for label, pages in items:
        _add_textbox(s, MARGIN_L, y, Inches(8), Inches(0.5),
                     label, font=BODY_FONT, size=13, color=WHITE)
        _add_textbox(s, Inches(10.5), y, Inches(2.5), Inches(0.5),
                     pages, font=BODY_FONT, size=13, color=WHITE)
        y += Inches(0.55)
    # Footer — match brand string used elsewhere
    _add_textbox(s, Inches(10.0), FOOTER_Y, Inches(3.0), Inches(0.25),
                 f"Autodecision     {page_num}",
                 font=BODY_FONT, size=8, color=RGBColor(0x9C, 0xA9, 0xB5),
                 align=PP_ALIGN.RIGHT)


def add_chapter_divider(prs, chapter_title, content_bullets, activities, page_num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, NAVY)
    _add_textbox(s, MARGIN_L, Inches(0.55), Inches(12), Inches(1.3),
                 chapter_title, font=TITLE_FONT, size=28, bold=True, color=WHITE)
    _add_rule(s, MARGIN_L, Inches(1.55), SLIDE_W - MARGIN_L - MARGIN_R,
              color=RGBColor(0x4A, 0x66, 0x7E), weight=0.5)

    _add_textbox(s, MARGIN_L, Inches(1.95), Inches(2.5), Inches(0.4),
                 "Section content", font=BODY_FONT, size=13, color=WHITE)
    body_tb = s.shapes.add_textbox(Inches(3.2), Inches(1.95), Inches(9.5), Inches(3))
    _add_paragraphs(body_tb.text_frame, content_bullets,
                    font=BODY_FONT, size=12, color=WHITE, bullet=True, leading=1.4)

    if activities:
        _add_textbox(s, MARGIN_L, Inches(4.6), Inches(2.5), Inches(0.4),
                     "Outputs", font=BODY_FONT, size=13, color=WHITE)
        out_tb = s.shapes.add_textbox(Inches(3.2), Inches(4.6), Inches(9.5), Inches(2.4))
        _add_paragraphs(out_tb.text_frame, activities,
                        font=BODY_FONT, size=12, color=WHITE, bullet=True, leading=1.4)

    _add_textbox(s, Inches(11.0), FOOTER_Y, Inches(2.0), Inches(0.25),
                 f"Autodecision Engine          {page_num}",
                 font=BODY_FONT, size=8, color=RGBColor(0x9C, 0xA9, 0xB5),
                 align=PP_ALIGN.RIGHT)


def _body_top(rule_y, *, subtitle=False):
    """Y where body content begins, given the rule_y returned by
    _add_action_title and an optional subtitle line."""
    return rule_y + (Inches(0.55) if subtitle else Inches(0.30))


def add_action_text_slide(prs, title, bullets, source, page_num, *,
                          prefix=None, subtitle=None):
    """Action title + bulleted body (used for Exec Summary etc)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    rule_y = _add_action_title(s, title, size=22, prefix=prefix)
    if subtitle:
        _add_textbox(s, MARGIN_L, rule_y + Inches(0.08),
                     SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.3),
                     subtitle, font=BODY_FONT, size=11,
                     color=GREY_FOOT, italic=True)
    body_y = _body_top(rule_y, subtitle=bool(subtitle))
    body_tb = s.shapes.add_textbox(MARGIN_L, body_y,
                                   SLIDE_W - MARGIN_L - MARGIN_R,
                                   FOOTER_Y - body_y - Inches(0.2))
    _add_paragraphs(body_tb.text_frame, bullets,
                    font=BODY_FONT, size=12, color=INK, bullet=True, leading=1.35)
    _add_footer(s, page_num, source=source)


def add_two_column_slide(prs, title, left_header, left_items, right_header,
                         right_items, source, page_num, *, prefix=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    rule_y = _add_action_title(s, title, size=22, prefix=prefix)
    body_y = _body_top(rule_y)

    col_w = (SLIDE_W - MARGIN_L - MARGIN_R - Inches(0.4)) / 2

    _add_textbox(s, MARGIN_L, body_y, col_w, Inches(0.4),
                 left_header, font=BODY_FONT, size=13, bold=True, color=MCK_BLUE)
    _add_rule(s, MARGIN_L, body_y + Inches(0.4), col_w, color=GREY_RULE)
    left_tb = s.shapes.add_textbox(MARGIN_L, body_y + Inches(0.55),
                                   col_w, Inches(5))
    _add_paragraphs(left_tb.text_frame, left_items, font=BODY_FONT, size=11,
                    color=INK, bullet=True, leading=1.35)

    rx = MARGIN_L + col_w + Inches(0.4)
    _add_textbox(s, rx, body_y, col_w, Inches(0.4),
                 right_header, font=BODY_FONT, size=13, bold=True,
                 color=ACCENT_CORAL)
    _add_rule(s, rx, body_y + Inches(0.4), col_w, color=GREY_RULE)
    right_tb = s.shapes.add_textbox(rx, body_y + Inches(0.55),
                                    col_w, Inches(5))
    _add_paragraphs(right_tb.text_frame, right_items, font=BODY_FONT, size=11,
                    color=INK, bullet=True, leading=1.35)

    _add_footer(s, page_num, source=source)


def add_three_column_grid(prs, title, columns, source, page_num, *, prefix=None):
    """columns = list of (header_label, accent_color, [bullets])"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    rule_y = _add_action_title(s, title, size=22, prefix=prefix)
    body_y = _body_top(rule_y)
    n = len(columns)
    gap = Inches(0.3)
    col_w = (SLIDE_W - MARGIN_L - MARGIN_R - gap * (n - 1)) / n

    for i, (header, color, items) in enumerate(columns):
        x = MARGIN_L + i * (col_w + gap)
        _add_textbox(s, x, body_y, col_w, Inches(0.4),
                     header, font=BODY_FONT, size=13, bold=True, color=color)
        _add_rule(s, x, body_y + Inches(0.4), col_w, color=color, weight=1.5)
        body = s.shapes.add_textbox(x, body_y + Inches(0.55), col_w, Inches(5))
        _add_paragraphs(body.text_frame, items, font=BODY_FONT, size=10,
                        color=INK, bullet=True, leading=1.3)

    _add_footer(s, page_num, source=source)


def add_table_slide(prs, title, headers, rows, col_widths, source, page_num,
                    *, subtitle=None, prefix=None, row_h=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    rule_y = _add_action_title(s, title, size=22, prefix=prefix)
    if subtitle:
        _add_textbox(s, MARGIN_L, rule_y + Inches(0.08),
                     SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.3),
                     subtitle, font=BODY_FONT, size=11,
                     color=GREY_FOOT, italic=True)
    table_top = _body_top(rule_y, subtitle=bool(subtitle))

    total_w = sum(col_widths)
    x = MARGIN_L
    for i, h in enumerate(headers):
        _add_textbox(s, x, table_top, col_widths[i], Inches(0.4), h,
                     font=BODY_FONT, size=11, bold=True, color=INK)
        x += col_widths[i]
    _add_rule(s, MARGIN_L, table_top + Inches(0.45), total_w,
              color=GREY_RULE, weight=1)

    # Auto-fit row height to available space if not specified, leaving
    # ample margin above the footer so source line never collides. Floor
    # is 0.28" so dense source tables (12-15 rows) still fit cleanly.
    if row_h is None:
        avail = FOOTER_Y - (table_top + Inches(0.55)) - Inches(0.45)
        per_row = avail / max(len(rows), 1)
        row_h = min(Inches(0.55), max(Inches(0.28), per_row))

    y = table_top + Inches(0.55)
    for row in rows:
        x = MARGIN_L
        for i, cell in enumerate(row):
            # Cell encoding:
            #   str                       → plain text in INK
            #   (str, color)              → colored text
            #   (str, color, "pill")      → status pill: colored disc + text
            kind = None
            cell_color = INK
            if isinstance(cell, tuple):
                if len(cell) == 3:
                    cell_text, cell_color, kind = cell
                else:
                    cell_text, cell_color = cell
            else:
                cell_text = cell
            if kind == "pill":
                _add_status_pill(s, x, y, col_widths[i], cell_text,
                                 color=cell_color)
            else:
                _add_textbox(s, x, y, col_widths[i], row_h, cell_text,
                             font=BODY_FONT, size=10, color=cell_color,
                             anchor=MSO_ANCHOR.TOP)
            x += col_widths[i]
        _add_rule(s, MARGIN_L, y + row_h, total_w,
                  color=GREY_RULE, weight=0.5)
        y += row_h

    _add_footer(s, page_num, source=source)


def add_recommendation_slide(prs, title, action, fields, source, page_num,
                             *, prefix=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    rule_y = _add_action_title(s, title, size=22, prefix=prefix)
    body_y = _body_top(rule_y)

    # Big call-out for action — no redundant label since the prefix above
    # already says RECOMMENDATION
    box_h = Inches(1.3)
    _add_filled_rect(s, MARGIN_L, body_y,
                     SLIDE_W - MARGIN_L - MARGIN_R, box_h, color=NAVY)
    # Bright accent stripe on left edge
    _add_filled_rect(s, MARGIN_L, body_y, Inches(0.08), box_h,
                     color=MCK_BLUE_BR)
    _add_textbox(s, MARGIN_L + Inches(0.4), body_y + Inches(0.22),
                 SLIDE_W - MARGIN_L - MARGIN_R - Inches(0.6), Inches(0.95),
                 action, font=TITLE_FONT, size=17, bold=True, color=WHITE)

    grid_top = body_y + box_h + Inches(0.25)
    cell_w = (SLIDE_W - MARGIN_L - MARGIN_R - Inches(0.4)) / 3
    cell_h = Inches(1.65)
    for idx, (label, value, color) in enumerate(fields):
        col = idx % 3
        row = idx // 3
        x = MARGIN_L + col * (cell_w + Inches(0.2))
        y = grid_top + row * (cell_h + Inches(0.15))
        _add_textbox(s, x, y, cell_w, Inches(0.3),
                     label, font=BODY_FONT, size=10, bold=True, color=color)
        _add_rule(s, x, y + Inches(0.32), cell_w, color=color, weight=1.2)
        _add_textbox(s, x, y + Inches(0.4), cell_w, cell_h - Inches(0.4),
                     value, font=BODY_FONT, size=9.5, color=INK)

    _add_footer(s, page_num, source=source)


# ---------------------------------------------------------------------------
# Chart helpers
# ---------------------------------------------------------------------------

def _mck_chart_setup():
    """Apply McKinsey-ish chart defaults to a fresh figure."""
    plt.rcParams.update({
        "font.family": ["Helvetica", "Arial", "sans-serif"],
        "font.size": 9,
        "axes.spines.top": False,
        "axes.spines.right": False,
        "axes.spines.left": False,
        "axes.edgecolor": "#666",
        "axes.labelcolor": "#1A1A1A",
        "xtick.color": "#1A1A1A",
        "ytick.color": "#1A1A1A",
        "axes.grid": True,
        "axes.grid.axis": "y",
        "grid.color": "#E9ECEF",
        "grid.linewidth": 0.5,
    })


def _make_council_mark(*, dpi=220):
    """The autodecision cover mark — 'Council to Synthesis'.

    Center synthesis node, inner pentagon of 5 persona dots (1st-order
    effects in each persona's color), and an outer ring of 2nd-order
    effects fanning out from each persona. Encodes the methodology
    visually: independent perspectives → 1st-order effects → 2nd-order
    cascade.

    Returns a transparent PNG buffer.
    """
    import numpy as np
    fig = plt.figure(figsize=(6, 6), dpi=dpi)
    ax = fig.add_axes([0, 0, 1, 1])

    # Persona palette — matches the radar chart elsewhere in the deck so
    # the cover and the council-depth slide read as one system.
    persona_colors = [
        ("#F39C12", "Optimist"),
        ("#C0392B", "Pessimist"),
        ("#9B59B6", "Competitor"),
        ("#054AD8", "Regulator"),
        ("#1F7A44", "Customer"),
    ]

    # Two rings: 1st-order (personas) and 2nd-order (effects cascade)
    n = 5
    r1 = 0.30
    r2 = 0.46
    cx, cy = 0.5, 0.5
    angles = [np.pi / 2 + 2 * np.pi * i / n for i in range(n)]
    points = [(cx + r1 * np.cos(a), cy + r1 * np.sin(a)) for a in angles]

    # Faint guide rings at both orbital radii
    theta = np.linspace(0, 2 * np.pi, 200)
    for r in (r1, r2):
        ax.plot(cx + r * np.cos(theta), cy + r * np.sin(theta),
                color=(0.05, 0.30, 0.85, 0.12),
                linewidth=0.5, linestyle=(0, (2, 5)), zorder=0)

    # 2nd-order effects: 3 small dots fanning from each persona's angle.
    # Drawn before convergence rays so the rays sit on top.
    fan = (-0.18, 0.0, 0.18)
    for i, a in enumerate(angles):
        color = persona_colors[i][0]
        x1, y1 = points[i]
        for da in fan:
            ang2 = a + da
            x2 = cx + r2 * np.cos(ang2)
            y2 = cy + r2 * np.sin(ang2)
            ax.plot([x1, x2], [y1, y2], color=color, alpha=0.32,
                    linewidth=0.9, solid_capstyle="round", zorder=1)
            ax.scatter([x2], [y2], s=70, color=color, alpha=0.70,
                       edgecolor="none", zorder=4)

    # Convergence rays — each persona to the synthesis center, gradient
    # from the persona's color (at outer end) toward navy (at center).
    # Render as many short segments to fake the gradient.
    SEGMENTS = 40
    for (px, py), (color, _) in zip(points, persona_colors):
        # Convert color hex → rgb 0-1
        cr = int(color[1:3], 16) / 255
        cg = int(color[3:5], 16) / 255
        cb = int(color[5:7], 16) / 255
        # Navy center color
        nr, ng, nb = 0x05 / 255, 0x1C / 255, 0x2C / 255
        for k in range(SEGMENTS):
            t0 = k / SEGMENTS
            t1 = (k + 1) / SEGMENTS
            x0 = px + (cx - px) * t0
            y0 = py + (cy - py) * t0
            x1 = px + (cx - px) * t1
            y1 = py + (cy - py) * t1
            # Color along the line: persona at t=0, navy at t=1
            t_mid = (t0 + t1) / 2
            r = cr + (nr - cr) * t_mid
            g = cg + (ng - cg) * t_mid
            b = cb + (nb - cb) * t_mid
            # Alpha is highest in the middle of the line, fading near ends
            a = 0.35 + 0.35 * np.sin(np.pi * t_mid)
            # Line tapers thinner near the center node
            lw = 1.6 * (1.0 - 0.55 * t_mid)
            ax.plot([x0, x1], [y0, y1],
                    color=(r, g, b, a), linewidth=lw, solid_capstyle="round",
                    zorder=2)

    # Persona dots — clean colored discs with white stroke
    for (px, py), (color, _) in zip(points, persona_colors):
        ax.scatter([px], [py], s=380, color=color,
                   edgecolor="white", linewidth=2.5, zorder=4)
        # Subtle outer halo for premium feel
        ax.scatter([px], [py], s=900, color=color,
                   alpha=0.10, edgecolor="none", zorder=3)

    # Synthesis center — larger navy disc with white stroke + accent ring
    ax.scatter([cx], [cy], s=1400, color="#051C2C",
               edgecolor="white", linewidth=3.0, zorder=6)
    # Inner accent dot in McKinsey-blue (subtle)
    ax.scatter([cx], [cy], s=180, color="#054AD8",
               edgecolor="none", zorder=7)
    # Faint navy halo
    ax.scatter([cx], [cy], s=2800, color="#051C2C",
               alpha=0.08, edgecolor="none", zorder=5)

    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")
    ax.set_aspect("equal")

    buf = io.BytesIO()
    fig.savefig(buf, format="png", transparent=True,
                bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    buf.seek(0)
    return buf


def _make_brand_mark_small(*, dpi=220):
    """Compact brand mark — center synthesis + 5 persona dots only.

    Drops the convergence rays, second-order ring, and halos so the
    mark stays legible at ~0.3 inch. Used for slide headers/footers
    next to the wordmark.

    Returns a transparent PNG buffer.
    """
    import numpy as np
    fig = plt.figure(figsize=(2, 2), dpi=dpi)
    ax = fig.add_axes([0, 0, 1, 1])

    persona_colors = ["#F39C12", "#C0392B", "#9B59B6", "#054AD8", "#1F7A44"]
    n = 5
    cx, cy = 0.5, 0.5
    r = 0.34
    angles = [np.pi / 2 + 2 * np.pi * i / n for i in range(n)]

    for color, a in zip(persona_colors, angles):
        x, y = cx + r * np.cos(a), cy + r * np.sin(a)
        ax.scatter([x], [y], s=140, color=color, edgecolor="none", zorder=2)

    ax.scatter([cx], [cy], s=320, color="#051C2C", edgecolor="none", zorder=3)

    ax.set_xlim(0.05, 0.95); ax.set_ylim(0.05, 0.95); ax.axis("off")
    ax.set_aspect("equal")

    buf = io.BytesIO()
    fig.savefig(buf, format="png", transparent=True,
                bbox_inches="tight", pad_inches=0.02)
    plt.close(fig)
    buf.seek(0)
    return buf


def make_effects_bar_chart(effects, *, w=8.5, h=5.0):
    """Horizontal bars: probability (with range), color by council_agreement."""
    _mck_chart_setup()
    fig, ax = plt.subplots(figsize=(w, h), dpi=200)
    labels = [e["label"] for e in effects]
    probs = [e["p"] for e in effects]
    agree = [e["agree"] for e in effects]
    y = list(range(len(labels)))
    # Color scale: 5 = navy, 4 = blue, 3 = bright blue, 2 = pale blue,
    # 1 = pale blue with grey outline (so the bar doesn't disappear into
    # the slide background for single-persona "specialist insight" effects).
    palette = {5: "#051C2C", 4: "#054AD8", 3: "#2196F3", 2: "#9DC3E6", 1: "#C9D6E2"}
    colors = [palette.get(a, "#9DC3E6") for a in agree]
    # 1-agreement bars get a subtle outline so they remain readable on
    # white background; higher-agreement bars use white edge (invisible).
    edge_colors = ["#7A8B9C" if a == 1 else "white" for a in agree]
    bars = ax.barh(y, probs, color=colors, edgecolor=edge_colors,
                   linewidth=0.8, height=0.62)
    # Probability range overlay (thin black line)
    for i, e in enumerate(effects):
        if e.get("range"):
            lo, hi = e["range"]
            ax.plot([lo, hi], [i, i], color="black", linewidth=1.1)
            ax.plot([lo, lo], [i - 0.18, i + 0.18], color="black", linewidth=1.1)
            ax.plot([hi, hi], [i - 0.18, i + 0.18], color="black", linewidth=1.1)
    # Probability labels at end of bar
    for i, p in enumerate(probs):
        ax.text(p + 0.015, i, f"{int(p*100)}%", va="center",
                fontsize=8.5, color="#1A1A1A")
    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=9)
    ax.invert_yaxis()
    ax.set_xlim(0, 1)
    ax.set_xticks([0, 0.25, 0.5, 0.75, 1.0])
    ax.set_xticklabels(["0%", "25%", "50%", "75%", "100%"], fontsize=8.5, color="#666")
    ax.tick_params(axis='y', length=0)
    ax.tick_params(axis='x', length=0)
    ax.set_axisbelow(True)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=200)
    plt.close(fig)
    buf.seek(0)
    return buf


def make_sensitivity_tornado(assumptions, *, w=8.5, h=5.0):
    """Tornado chart: bar length = effects_impacted * sensitivity_weight."""
    _mck_chart_setup()
    fig, ax = plt.subplots(figsize=(w, h), dpi=200)
    sens_w = {"HIGH": 1.0, "MEDIUM": 0.55, "LOW": 0.2}
    weighted = [(a["label"], a["effects"] * sens_w.get(a["sensitivity"], 0.5),
                 a["sensitivity"]) for a in assumptions]
    weighted.sort(key=lambda x: x[1])
    labels = [x[0] for x in weighted]
    vals = [x[1] for x in weighted]
    sens = [x[2] for x in weighted]
    palette = {"HIGH": "#054AD8", "MEDIUM": "#2196F3", "LOW": "#9DC3E6"}
    colors = [palette[s] for s in sens]
    y = list(range(len(labels)))
    ax.barh(y, vals, color=colors, edgecolor="white", height=0.62)
    for i, (v, s_lvl) in enumerate(zip(vals, sens)):
        ax.text(v + 0.05, i, s_lvl, va="center", fontsize=8.5,
                color="#1A1A1A")
    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=9)
    ax.set_xticks([])
    ax.tick_params(axis='y', length=0)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=200)
    plt.close(fig)
    buf.seek(0)
    return buf


def add_section_divider(prs, section_num, section_title, lead_question,
                         page_num):
    """Dark-navy chapter divider with big serif section title and a leading
    question — anchors deck structure between major parts."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, NAVY)

    # Big section number top-left, in McKinsey blue
    _add_textbox(s, MARGIN_L, Inches(0.55), Inches(2), Inches(0.6),
                 f"PART {section_num}", font=BODY_FONT, size=12, bold=True,
                 color=MCK_BLUE_BR)

    # Section title — large serif white. Auto-shrink for multi-line titles
    # so a 2-line "How we explored this decision" doesn't feel oversized.
    is_multiline = "\n" in section_title or len(section_title) > 24
    title_size = 36 if is_multiline else 44
    title_tb = s.shapes.add_textbox(MARGIN_L, Inches(2.55),
                                     Inches(11), Inches(2.2))
    tf = title_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = section_title
    run.font.name = TITLE_FONT
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p.line_spacing = 1.05

    # Thin rule below title in pale blue
    _add_rule(s, MARGIN_L, Inches(4.3), Inches(8.0),
              color=RGBColor(0x4A, 0x66, 0x7E), weight=1)

    # Leading question — italic, smaller, like a chapter epigraph
    _add_textbox(s, MARGIN_L, Inches(4.5), Inches(11.5), Inches(1.5),
                 lead_question, font=BODY_FONT, size=14,
                 color=RGBColor(0xB8, 0xC4, 0xD0), italic=True)

    # Footer page number in subtle grey
    _add_textbox(s, Inches(10.0), FOOTER_Y, Inches(3.0), Inches(0.25),
                 f"Autodecision     {page_num}",
                 font=BODY_FONT, size=8, color=RGBColor(0x9C, 0xA9, 0xB5),
                 align=PP_ALIGN.RIGHT)


def add_matrix_2x2_slide(prs, title, axes, items, source, page_num, *,
                         prefix=None, quadrant_labels=None):
    """2x2 strategic matrix — bubbles plotted on (x, y, size) tuples.
    axes: dict with keys x_low, x_high, y_low, y_high, x_label, y_label.
    items: list of dicts {label, x, y, size, color, note (optional)}.
    quadrant_labels: dict {'tl', 'tr', 'bl', 'br'} → labels (optional)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    rule_y = _add_action_title(s, title, size=22, prefix=prefix)
    body_y = _body_top(rule_y)

    fig, ax = plt.subplots(figsize=(10, 5.0), dpi=200)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    # Quadrant background tints
    ax.add_patch(Rectangle((0, 0.5), 0.5, 0.5, facecolor="#FCE8EC",
                           edgecolor="none", zorder=0))  # top-left worst
    ax.add_patch(Rectangle((0.5, 0.5), 0.5, 0.5, facecolor="#FFF4D6",
                           edgecolor="none", zorder=0))
    ax.add_patch(Rectangle((0, 0), 0.5, 0.5, facecolor="#F2F4F7",
                           edgecolor="none", zorder=0))
    ax.add_patch(Rectangle((0.5, 0), 0.5, 0.5, facecolor="#E3F0E1",
                           edgecolor="none", zorder=0))
    # Quadrant labels
    if quadrant_labels:
        ax.text(0.02, 0.97, quadrant_labels.get("tl", ""), va="top",
                fontsize=9, fontweight="bold", color="#A33B3B")
        ax.text(0.98, 0.97, quadrant_labels.get("tr", ""), va="top", ha="right",
                fontsize=9, fontweight="bold", color="#8B6F1A")
        ax.text(0.02, 0.03, quadrant_labels.get("bl", ""), va="bottom",
                fontsize=9, fontweight="bold", color="#666")
        ax.text(0.98, 0.03, quadrant_labels.get("br", ""), va="bottom", ha="right",
                fontsize=9, fontweight="bold", color="#1F7A44")
    # Axis lines
    ax.axhline(0.5, color="#666", linewidth=0.6)
    ax.axvline(0.5, color="#666", linewidth=0.6)
    # Bubbles — label_pos controls placement to avoid stacking when bubbles
    # cluster. Values: "below" (default), "above", "left", "right".
    POS_OFFSETS = {
        "below":  (0,    -0.055, "center", "top"),
        "above":  (0,     0.055, "center", "bottom"),
        "left":   (-0.04, 0,     "right",  "center"),
        "right":  (0.04,  0,     "left",   "center"),
    }
    for it in items:
        # Guardrail: never plot a bubble exactly on the quadrant dividers
        # (x=0.5 or y=0.5). A bubble on the line reads as ambiguous between
        # two quadrants, defeating the matrix's purpose. Nudge by 0.05 in
        # whichever direction the bubble was leaning.
        if abs(it["x"] - 0.5) < 0.04:
            it = dict(it, x=0.45 if it["x"] < 0.5 else 0.55)
        if abs(it["y"] - 0.5) < 0.04:
            it = dict(it, y=0.45 if it["y"] < 0.5 else 0.55)
        c = it.get("color", "#054AD8")
        ax.scatter(it["x"], it["y"], s=it["size"] * 25, color=c,
                   edgecolor="white", linewidth=1.4, zorder=3, alpha=0.85)
        # bubble size affects label offset so labels don't overlap dots
        radius_offset = 0.015 + (it["size"] / 16) * 0.025
        dx, dy, ha, va = POS_OFFSETS.get(it.get("label_pos", "below"),
                                          POS_OFFSETS["below"])
        # scale offset by direction
        if dy < 0: dy -= radius_offset
        if dy > 0: dy += radius_offset
        if dx < 0: dx -= radius_offset
        if dx > 0: dx += radius_offset
        ax.text(it["x"] + dx, it["y"] + dy, it["label"], fontsize=8.5,
                ha=ha, va=va, color="#1A1A1A", zorder=4,
                fontweight="bold")
        if it.get("note"):
            # note always opposite of label
            note_dy = +0.06 if dy <= 0 else -0.06
            note_va = "bottom" if note_dy > 0 else "top"
            ax.text(it["x"], it["y"] + note_dy, it["note"], fontsize=7,
                    ha="center", va=note_va, color="#555", zorder=4,
                    style="italic")
    # Axis labels
    ax.set_xlabel(axes["x_label"], fontsize=10, color="#1A1A1A",
                  fontweight="bold", labelpad=10)
    ax.set_ylabel(axes["y_label"], fontsize=10, color="#1A1A1A",
                  fontweight="bold", labelpad=10)
    ax.set_xticks([0, 0.5, 1]); ax.set_yticks([0, 0.5, 1])
    ax.set_xticklabels([axes["x_low"], "", axes["x_high"]],
                        fontsize=8, color="#666")
    ax.set_yticklabels([axes["y_low"], "", axes["y_high"]],
                        fontsize=8, color="#666")
    for spine in ax.spines.values():
        spine.set_color("#666")
        spine.set_linewidth(0.6)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=200)
    plt.close(fig)
    buf.seek(0)
    s.shapes.add_picture(buf, MARGIN_L, body_y, width=Inches(8.5))

    # Legend column
    com_x = MARGIN_L + Inches(8.7)
    com_w = SLIDE_W - com_x - MARGIN_R
    _add_textbox(s, com_x, body_y, com_w, Inches(0.35),
                 "READ", font=BODY_FONT, size=9.5, bold=True, color=MCK_BLUE)
    _add_rule(s, com_x, body_y + Inches(0.34), com_w,
              color=MCK_BLUE, weight=1.2)
    legend = []
    for it in items:
        legend.append(f"{it['label']} — {it.get('legend', '')}")
    legend_tb = s.shapes.add_textbox(com_x, body_y + Inches(0.45),
                                     com_w, Inches(5))
    _add_paragraphs(legend_tb.text_frame, legend, font=BODY_FONT, size=9.5,
                    color=INK, bullet=True, leading=1.4)

    _add_footer(s, page_num, source=source)


def add_radar_slide(prs, title, axes, personas, source, page_num, *,
                    prefix=None, commentary_lines=None):
    """Per-persona radar across N axes. personas = list of dicts
    {name, scores: [n], color}.

    Note: keep axis labels SHORT (single word ideally) — polar plots in
    matplotlib place tick labels external to the circle, so long labels
    push off the figure edge."""
    import numpy as np
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    rule_y = _add_action_title(s, title, size=22, prefix=prefix)
    body_y = _body_top(rule_y)

    n = len(axes)
    angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
    angles += [angles[0]]

    # Wider figure with extra margin so external labels fit
    fig = plt.figure(figsize=(6.0, 4.6), dpi=200)
    ax = fig.add_axes([0.18, 0.10, 0.64, 0.80], polar=True)

    for p in personas:
        scores = list(p["scores"]) + [p["scores"][0]]
        ax.plot(angles, scores, color=p["color"], linewidth=1.6,
                label=p["name"])
        ax.fill(angles, scores, color=p["color"], alpha=0.08)

    # Axes config
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(axes, fontsize=8, color="#1A1A1A")
    ax.set_ylim(0, 5)
    ax.set_yticks([1, 2, 3, 4, 5])
    ax.set_yticklabels(["", "", "", "", ""])
    ax.tick_params(axis='x', pad=4)
    ax.grid(color="#D4D9DD", linewidth=0.5)
    ax.spines['polar'].set_color("#D4D9DD")
    ax.set_facecolor("#FAFBFC")
    # Legend below the chart (horizontal) so it doesn't clip axis labels
    ax.legend(loc="lower center", bbox_to_anchor=(0.5, -0.15), ncol=5,
              fontsize=7.5, frameon=False, columnspacing=1.0,
              handletextpad=0.4)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=200,
                pad_inches=0.15)
    plt.close(fig)
    buf.seek(0)
    # Available height in body, capped so chart never drops below footer
    chart_h = min(Inches(5.0), FOOTER_Y - body_y - Inches(0.3))
    s.shapes.add_picture(buf, MARGIN_L, body_y,
                         width=Inches(6.5), height=chart_h)

    # Commentary right side
    if commentary_lines:
        com_x = MARGIN_L + Inches(6.8)
        com_w = SLIDE_W - com_x - MARGIN_R
        body = s.shapes.add_textbox(com_x, body_y, com_w, Inches(5))
        _add_paragraphs(body.text_frame, commentary_lines,
                        font=BODY_FONT, size=10, color=INK,
                        bullet=True, leading=1.4)

    _add_footer(s, page_num, source=source)


def add_chart_slide(prs, title, chart_buf, commentary_lines, source, page_num,
                    *, chart_w=Inches(8.0), chart_x=None, prefix=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    rule_y = _add_action_title(s, title, size=22, prefix=prefix)
    body_y = _body_top(rule_y)
    cx = chart_x if chart_x else MARGIN_L
    s.shapes.add_picture(chart_buf, cx, body_y, width=chart_w)
    if commentary_lines:
        com_x = cx + chart_w + Inches(0.3)
        com_w = SLIDE_W - com_x - MARGIN_R
        body = s.shapes.add_textbox(com_x, body_y, com_w, Inches(5.2))
        _add_paragraphs(body.text_frame, commentary_lines,
                        font=BODY_FONT, size=10, color=INK, bullet=True, leading=1.4)
    _add_footer(s, page_num, source=source)


# ---------------------------------------------------------------------------
# Adobe deck content (hardcoded)
# ---------------------------------------------------------------------------

def build_adobe_deck(out_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- 1. Cover
    add_title_cover(
        prs,
        title="Should Adobe deprecate\nPhotoshop's UI within 3 years?",
        subtitle="Decision Brief — possibility map and recommendation",
        footer_org="Autodecision",
        date="April 2026",
    )

    # --- 2. TOC — section labels stay clean (no "Part X" prefix); slide
    # ranges, not just start numbers
    add_toc_dark(prs, [
        ("Recommendation",                "3"),
        ("How we explored this decision", "4 – 7"),
        ("What the council found",        "8 – 13"),
        ("Path forward",                  "14 – 15"),
        ("Appendix — Sources",            "16"),
    ], page_num=2)

    # --- 3. Executive summary / recommendation up front (pyramid principle)
    add_recommendation_slide(
        prs,
        prefix="1/  RECOMMENDATION",
        title="The aggressive deprecation path survives <15% of plausible scenarios; ship a segmented dual-SKU instead",
        action="Do NOT deprecate the Photoshop UI. Ship Photoshop Pro (UI + 10-year LTS) alongside Photoshop Agentic (Firefly AI Assistant native), with binding public reversal triggers.",
        fields=[
            ("CONFIDENCE", "HIGH on do-not-deprecate · MEDIUM-HIGH on segmented+LTS structure · MEDIUM on board willingness to publish reversal triggers", MCK_BLUE),
            ("DEPENDS ON", "Board willing to publish binding reversal triggers · LTS commitment credible to pros · Segments are separable · Adobe execution amortizes dual-SKU R&D", MCK_BLUE_BR),
            ("MONITOR", "Firefly AI Assistant beta NPS · Acrobat UI backlash trajectory · OpenAI/Anthropic creative product launches · Figma post-IPO M&A · Enterprise ELA renewal cohort", ACCENT_CORAL),
            ("PRE-MORTEM", "h6 without binding triggers (P=0.30) · Foundation model plateau (P=0.25) · Pro community rejects LTS as marketing-only (P=0.20) · Internal resourcing split degrades both SKUs (P=0.55)", MCK_BLUE_BR),
            ("DEEPEST DISAGREEMENT", "Optimist: 55% first-mover platform lock compounds. Pessimist: 75% reversal-impossible franchise damage. Disagreement IS the data.", MCK_BLUE),
            ("REVIEW TRIGGER", "Firefly AI Assistant beta NPS <40 · OpenAI/Anthropic ships native creative · Adobe board refuses public reversal triggers · Q2 FY26 Digital Media growth <10%", ACCENT_CORAL),
        ],
        source="Source: Decision Brief (full mode, 2 iterations, council 5/5 endorsement on h6) · [G15], [C3:council], [C4:pessimist], [C17:council]",
        page_num=3,
    )

    # --- 4. Section divider: how we explored
    add_section_divider(
        prs,
        section_num="2",
        section_title="How we explored\nthis decision",
        lead_question="A 5-persona council, 6 competing hypotheses, two iterations to convergence, and 12 adversarial stress scenarios — what survived?",
        page_num=4,
    )

    # --- 5. Methodology
    add_action_text_slide(
        prs,
        prefix="2.1/  HOW WE ANALYZED",
        title="The analysis explored 6 hypotheses across 2 iterations of a 5-persona council, stress-tested by 12 adversarial scenarios",
        bullets=[
            "5 personas spawned as separate subagents — Optimist, Pessimist, Competitor, Regulator, Customer — each with independent context",
            "6 hypotheses generated, ranging from aggressive deprecation (h1) to additive Adobe-actual strategy (h3) to the synthesis winner (h6 segmented dual-SKU + LTS + binding reversal triggers)",
            "2 iterations — convergence reached at iteration 2 with 78% assumption stability, 1 productive contradiction remaining",
            "Adversary phase generated 12 stress scenarios across worst-cases, black swans, and irrational actors",
            "Sensitivity analysis on 10 key assumptions — 3 high-sensitivity, 5 medium, 2 low",
            "Real-time anchor: Adobe announced Firefly AI Assistant as additive on April 15, 2026 [G15] — three days before this analysis. The announced strategy is h3, the recommendation upgrades to h6.",
        ],
        source="Source: Decision Brief Convergence Log; engine-protocol.md",
        page_num=5,
    )

    # --- 6. Hypotheses tracker
    add_table_slide(
        prs,
        prefix="2.2/  HYPOTHESES",
        title="Six hypotheses surfaced; only h6 survived 5/5 council endorsement and adversarial stress",
        headers=["#", "Hypothesis", "Status", "Key assumptions"],
        rows=[
            ["1",
             "Aggressive deprecation wins (3-year UI sunset)",
             ("Conditional · joint P 12-20%", RGBColor(0xC0, 0x39, 0x2B), "pill"),
             "Agentic is terminal paradigm · 36-month parity · Foundation labs enter creative"],
            ["2",
             "User revolt destroys franchise",
             ("Risk path of h1", RGBColor(0xC0, 0x39, 0x2B), "pill"),
             "Muscle-memory switching cost real · Narrative damage sticky · Enterprise procurement revolts"],
            ["3",
             "Hybrid additive dominates (Adobe's actual April 2026 strategy [G15])",
             ("Strongly supported", MCK_BLUE, "pill"),
             "Both modes coexist · Market rewards additive · AI Assistant quality holds"],
            ["4",
             "Competitive window closes by 2028",
             ("Conditional on foundation lab entry", GREY_FOOT, "pill"),
             "Foundation labs enter creative · Category narrative matters · Translates to valuation"],
            ["5",
             "Segmentation is the right structure",
             ("Conditional on segment separability", GREY_FOOT, "pill"),
             "Segments separable · Dual tier doesn't confuse · Pricing power preserved"],
            ["6",
             "Segmented + 10yr LTS + binding reversal triggers",
             ("Leading · 5/5 council", RGBColor(0x1F, 0x7A, 0x44), "pill"),
             "Dual SKU separable · LTS credible · Reversal triggers enforceable"],
        ],
        col_widths=[Inches(0.5), Inches(4.5), Inches(2.7), Inches(4.6)],
        source="Source: Decision Brief — Hypotheses Explored",
        page_num=6,
    )

    # --- 7. NEW: 2x2 strategic positioning matrix
    add_matrix_2x2_slide(
        prs,
        prefix="2.3/  STRATEGIC POSITIONING",
        title="Plotted on reversibility × franchise risk, h6 sits in the manageable quadrant; h1 sits where every adversarial scenario amplifies",
        axes={
            "x_label": "Reversibility of decision",
            "y_label": "Franchise risk if wrong",
            "x_low": "Low (one-way door)",
            "x_high": "High (recoverable)",
            "y_low": "Contained",
            "y_high": "Catastrophic",
        },
        quadrant_labels={
            "tl": "AVOID",
            "tr": "MANAGE WITH RIGOR",
            "bl": "MONITOR",
            "br": "DEFAULT",
        },
        items=[
            {"label": "h1", "x": 0.10, "y": 0.85, "size": 8,
             "color": "#C0392B", "label_pos": "right",
             "legend": "Aggressive deprecation — irreversible, catastrophic if wrong (joint P 12-20%)"},
            {"label": "h2", "x": 0.20, "y": 0.72, "size": 6,
             "color": "#E08673", "label_pos": "right",
             "legend": "User revolt — risk path of h1, same quadrant"},
            {"label": "h4", "x": 0.40, "y": 0.55, "size": 7,
             "color": "#9CA3AF", "label_pos": "right",
             "legend": "Window-closes — conditional, indeterminate position"},
            {"label": "h5", "x": 0.62, "y": 0.45, "size": 9,
             "color": "#2196F3", "label_pos": "left",
             "legend": "Segmentation — direction right, structure incomplete"},
            {"label": "h3", "x": 0.83, "y": 0.20, "size": 12,
             "color": "#1F7A44", "label_pos": "below",
             "legend": "Hybrid additive — Adobe's actual April 2026 strategy [G15]"},
            {"label": "h6", "x": 0.93, "y": 0.35, "size": 16,
             "color": "#054AD8", "label_pos": "left",
             "note": "RECOMMENDED",
             "legend": "Segmented + LTS + binding reversal triggers — captures growth while preserving optionality"},
        ],
        source="Source: Synthesis from Decision Brief — bubble size = council_agreement × probability of dominance",
        page_num=7,
    )

    # --- 8. Section divider: what the council found
    add_section_divider(
        prs,
        section_num="3",
        section_title="What the council\nfound",
        lead_question="Where the personas converged, where they diverged, and what the disagreement itself reveals about the decision.",
        page_num=8,
    )

    # --- 9. Effects map (high-confidence) bar chart
    high_conf_effects = [
        {"label": "Pro segment churn 18-28% IF aggressive deprecation",
         "p": 0.70, "agree": 5, "range": (0.55, 0.85)},
        {"label": "Installed base retained under additive (h3)",
         "p": 0.80, "agree": 5, "range": (0.70, 0.90)},
        {"label": "Short-term revenue trough (12% → 4-7% YoY)",
         "p": 0.70, "agree": 4, "range": (0.55, 0.80)},
        {"label": "Reversal impossible after announcement",
         "p": 0.75, "agree": 3, "range": (0.65, 0.85)},
        {"label": "Stock compression 25-45% on aggressive announcement",
         "p": 0.65, "agree": 4, "range": (0.50, 0.78)},
        {"label": "AI Assistant captures 80M+ freemium MAUs at 2-3x",
         "p": 0.65, "agree": 4, "range": (0.50, 0.75)},
        {"label": "Execution risk: 30-yr feature surface in 36 months",
         "p": 0.65, "agree": 4, "range": (0.55, 0.78)},
        {"label": "h6 binding reversal triggers eliminate one-way-door risk",
         "p": 0.65, "agree": 4, "range": (0.55, 0.75)},
        {"label": "Acrobat UI precedent escalates to Photoshop",
         "p": 0.65, "agree": 3, "range": (0.50, 0.78)},
        {"label": "Firefly flywheel compounds via session data",
         "p": 0.60, "agree": 3, "range": (0.45, 0.75)},
        {"label": "Platform moat deepens under additive",
         "p": 0.60, "agree": 4, "range": (0.50, 0.70)},
        {"label": "h6 dual SKU separates cleanly (5/5 council)",
         "p": 0.55, "agree": 5, "range": (0.45, 0.65)},
    ]
    add_chart_slide(
        prs,
        prefix="3.1/  EFFECTS MAP",
        title="The strongest signals point to franchise risk under aggressive paths and quiet compounding under additive",
        chart_buf=make_effects_bar_chart(high_conf_effects, w=8.0, h=4.6),
        commentary_lines=[
            "Top 12 effects ranked by council_agreement × probability",
            "Bar color = council agreement (navy = 5/5, fading to pale = 1/5)",
            "Black bracket = persona disagreement range — width is the uncertainty signal",
            "Two of the three highest-agreement effects (5/5) are downside under aggressive paths",
            "h6 effects (last row) emerge in iteration 2 with full council endorsement",
        ],
        source="Source: Decision Brief — Effects Map (top 15 by council_agreement × probability) · [C7:council], [C10:council], [C16:council]",
        page_num=9,
        chart_w=Inches(8.0),
    )

    # --- 10. Council dynamics
    add_two_column_slide(
        prs,
        prefix="3.2/  COUNCIL DYNAMICS",
        title="The deepest disagreement — competitive urgency vs franchise preservation — IS the data",
        left_header="STRONGEST: REGULATOR (#1 PEER RANK)",
        left_items=[
            "Won on structural depth — EU AI Act obligations, Firefly training-data suit risk, fiduciary / Delaware business judgment rule, ADA Section 508 implications of chat-only UI",
            "Surfaced TWO of the alternatives that synthesized into h6: dual-SKU LTS + binding reversal triggers",
            "Identified h6 as protected under business judgment rule — reversibility insulates the board even if strategy fails",
            "Compliance moat (AI Act + copyright + enterprise DPA) buys 18-24 months regardless of UI strategy",
        ],
        right_header="DEEPEST DIVIDE: URGENCY VS PRESERVATION",
        right_items=[
            "Optimist + Competitor: narrow window, first-mover lock compounds (P=0.55 [C2:optimist])",
            "Pessimist + Regulator + Customer: catastrophic franchise downside under forced transition (P=0.75 [C4:pessimist] reversal impossible)",
            "Disagreement range on foundation_labs_enter_creative: 0.30-0.70 [C3:council] — the disagreement IS the data",
            "Optimist was weakest (#5 peer rank) — joint-probability inflation flagged by 4/4 reviewers",
            "BLIND SPOT caught: 4/5 personas missed agentic inference COGS — ~10-50× per session, material at 41M subs",
        ],
        source="Source: Decision Brief — Council Dynamics · [C2:optimist], [C2:pessimist], [C3:council]",
        page_num=10,
    )

    # --- 11. NEW: Per-persona radar
    add_radar_slide(
        prs,
        prefix="3.3/  COUNCIL DEPTH",
        title="Regulator scored highest on structural depth and counter-factual rigor; Optimist's narrow ranges drove its #5 peer rank",
        axes=[
            "Effects",
            "Range width",
            "Adversarial",
            "Counter-factual",
            "Quant. rigor",
        ],
        personas=[
            {"name": "Optimist",   "scores": [3, 2, 2, 1, 3], "color": "#F39C12"},
            {"name": "Pessimist",  "scores": [4, 4, 5, 4, 4], "color": "#C0392B"},
            {"name": "Competitor", "scores": [3, 3, 4, 3, 3], "color": "#9B59B6"},
            {"name": "Regulator",  "scores": [4, 4, 4, 5, 5], "color": "#054AD8"},
            {"name": "Customer",   "scores": [3, 3, 2, 3, 3], "color": "#1F7A44"},
        ],
        commentary_lines=[
            "Each persona scored 1-5 across five rigor dimensions, derived from peer-review and critique outputs",
            "Regulator's structural depth surfaced both alternatives that synthesized into h6 — dual-SKU LTS and binding reversal triggers",
            "Optimist's narrow range and weak counter-factual depth flagged in 4/4 reviews — joint-probability inflation pattern",
            "Pessimist + Regulator together cover the full adversarial surface; their alignment on irreversibility risk drove the h6 design",
            "Customer's lower adversarial coverage is structural — the role focuses on adoption value, not threats",
        ],
        source="Source: Synthesized from Decision Brief — Council Dynamics + Critique outputs",
        page_num=11,
    )

    # --- 12. Section divider: stress test
    # (Adversarial + sensitivity sit under Part 3 too; no extra divider — keep
    #  to 3 dividers to avoid over-segmenting a 16-slide deck.)

    # --- 12. Adversarial scenarios
    add_three_column_grid(
        prs,
        prefix="3.4/  ADVERSARIAL PRESSURE",
        title="Twelve adversarial scenarios — eight push toward kill-aggressive or segmented; only h6's reversal triggers neutralize all of them",
        columns=[
            ("WORST CASES", MCK_BLUE, [
                "Foundation model capabilities plateau by 2027 (P=0.25)",
                "Pro community organizes mass migration via plugin APIs (P=0.20)",
                "EU AI Act enforcement hits Firefly mid-transition (P=0.20)",
                "Senior ML talent departs during stock-compression trough (P=0.30)",
            ]),
            ("BLACK SWANS", ACCENT_CORAL, [
                "Apple ships Vision Pro agentic Final Cut + Photos Dec 2027 (P=0.15)",
                "Figma acquires Runway or Midjourney with post-IPO war chest (P=0.15)",
                "CEO Narayen retires mid-transition; successor reverses (P=0.10)",
                "China bans Creative Cloud during trade escalation; $2-3B revenue (P=0.10)",
            ]),
            ("IRRATIONAL ACTORS", MCK_BLUE_BR, [
                "F500 flagship client publicly abandons Adobe; narrative contagion (P=0.15)",
                "Viral 'Adobe killed Photoshop' video from top creator (P=0.20)",
                "Activist investor takes 5% stake demanding reversal during trough (P=0.15)",
                "Board member leaks internal concerns; stock further compresses (P=0.08)",
            ]),
        ],
        source="Source: Decision Brief — Adversarial Scenarios · [C23:adversary] through [C35:adversary]",
        page_num=12,
    )

    # --- 13. Sensitivity tornado
    assumptions = [
        {"label": "Adobe ships agentic parity in 36 months",
         "sensitivity": "HIGH", "effects": 3},
        {"label": "Agentic is terminal UI paradigm for creative",
         "sensitivity": "HIGH", "effects": 3},
        {"label": "Foundation labs enter creative by 2028",
         "sensitivity": "HIGH", "effects": 2},
        {"label": "Adobe board willing to publish binding reversal triggers",
         "sensitivity": "HIGH", "effects": 2},
        {"label": "Market rewards additive strategy for cash-cow incumbents",
         "sensitivity": "MEDIUM", "effects": 3},
        {"label": "Segments are separable (pros vs prosumers)",
         "sensitivity": "MEDIUM", "effects": 3},
        {"label": "Enterprise procurement revolts against forced UI transition",
         "sensitivity": "MEDIUM", "effects": 2},
        {"label": "Muscle-memory switching cost real for 20+ year pros",
         "sensitivity": "MEDIUM", "effects": 2},
        {"label": "Narrative damage from deprecation announcement is sticky",
         "sensitivity": "MEDIUM", "effects": 2},
        {"label": "LTS commitment on Pro SKU is credible (not marketing-only)",
         "sensitivity": "MEDIUM", "effects": 1},
    ]
    add_chart_slide(
        prs,
        prefix="3.5/  SENSITIVITY",
        title="Three high-sensitivity assumptions are load-bearing — if any one fails, aggressive deprecation becomes indefensible",
        chart_buf=make_sensitivity_tornado(assumptions, w=8.0, h=4.6),
        commentary_lines=[
            "Bar length = sensitivity weight × number of effects impacted",
            "Dark-blue bar = HIGH sensitivity, lighter = MEDIUM",
            "All four HIGH-sensitivity assumptions are external — Adobe controls none directly",
            "h6's binding reversal triggers convert assumption-failure from catastrophic to recoverable",
            "If 36-month parity slips (>60% probability per [C1:council]), aggressive deprecation guarantees the partial-agentic awkward state — exactly what h6 is structured to prevent",
        ],
        source="Source: Decision Brief — Key Assumptions · sensitivity per effects-impacted × HIGH/MEDIUM/LOW weight",
        page_num=13,
        chart_w=Inches(8.0),
    )

    # --- 14. Section divider: path forward
    add_section_divider(
        prs,
        section_num="4",
        section_title="Path forward",
        lead_question="Six checkpoints with binding reversal triggers — every gate has a kill criterion before the next commitment.",
        page_num=14,
    )

    # --- 15. Decision timeline
    add_table_slide(
        prs,
        prefix="4.1/  ROADMAP",
        title="The recommendation maps to a six-checkpoint roadmap with clear kill criteria at each stage",
        subtitle="Every checkpoint has a binding reversal trigger; failure at any gate triggers either reversal or extension of LTS",
        headers=["When", "Action", "Decision point", "Kill criteria"],
        rows=[
            ["Q2 FY26",
             "Firefly AI Assistant beta launches",
             "Beta NPS ≥ 50 from pros",
             ("Beta NPS < 30 → reassess", RGBColor(0xC0, 0x39, 0x2B))],
            ["Q3 FY26",
             "Publish Pro LTS commitment + Agentic SKU roadmap",
             "Public commitment with binding reversal triggers",
             ("Board refuses to publish → fall back to pure h3", RGBColor(0xC0, 0x39, 0x2B))],
            ["Q4 FY26",
             "Firefly AI Assistant GA across CC apps",
             "GA launches on schedule",
             ("Ship slip > 2 quarters → narrative damage risk", RGBColor(0xC0, 0x39, 0x2B))],
            ["FY27",
             "Agentic SKU builds; Pro SKU continues UI investment",
             "Agentic ARPU $15-20/mo; Pro maintained",
             ("Pro churn > 5% OR Agentic ARPU < $10 → reversal", RGBColor(0xC0, 0x39, 0x2B))],
            ["FY28",
             "Foundation lab entry decision point",
             "OpenAI/Anthropic creative products shipped or not",
             ("Foundation lab dominant category → accelerate", GREY_FOOT)],
            ["FY29",
             "3-year review point (original deprecation date)",
             "Reaffirm LTS; evaluate market shift",
             ("Pro churn > 8% → reversal; <3% + strong agentic → continue", RGBColor(0x1F, 0x7A, 0x44))],
        ],
        col_widths=[Inches(0.9), Inches(3.5), Inches(3.6), Inches(4.3)],
        source="Source: Decision Brief — Appendix A: Decision Timeline · binding reversal triggers per [C17:council]",
        page_num=15,
    )

    # --- 16. Sources (appendix)
    add_table_slide(
        prs,
        prefix="APPENDIX  SOURCES",
        title="Selected sources",
        subtitle="Tags appear inline throughout the deck — [G#] = ground data, [C#:persona] = council finding, [C#:adversary] = adversarial scenario",
        headers=["Tag", "Type", "Claim", "Source"],
        rows=[
            ["G2",  "Ground",   "Creative Cloud: 41M paid subs YE 2025; 80M+ freemium MAUs",                            "prodesigntools.com, qz.com"],
            ["G3",  "Ground",   "Adobe Acrobat UI redesign backlash sustained and unresolved",                          "acrobat.uservoice.com"],
            ["G15", "Ground",   "Firefly AI Assistant announced April 15 2026 (additive, not replacement)",             "blog.adobe.com, news.adobe.com"],
            ["G17", "Ground",   "Firefly AI Assistant supports third-party models (Anthropic Claude)",                  "appleinsider.com"],
            ["C3",  "Council",  "Foundation labs enter creative by 2028 P=0.5 joint",                                   "iter-1 council"],
            ["C4",  "Council",  "Reversal impossible P=0.75 after announcement",                                        "iter-1 pessimist"],
            ["C7",  "Council",  "Pro churn spike P=0.70 with 18-28% magnitude",                                         "iter-1 council"],
            ["C10", "Council",  "Installed base retained under additive P=0.80",                                        "iter-1 council"],
            ["C17", "Council",  "Binding reversal triggers at 8% pro churn",                                            "iter-2 council"],
            ["C23", "Adversary","Foundation model plateau P=0.25 (worst case)",                                         "iter-1 adversary"],
            ["C32", "Adversary","F500 flagship client abandons Adobe P=0.15 (irrational actor)",                        "iter-1 adversary"],
        ],
        col_widths=[Inches(0.7), Inches(1.2), Inches(7.5), Inches(2.9)],
        source="Source: Decision Brief — Sources (full list in companion brief markdown)",
        page_num=16,
    )

    prs.save(out_path)
    print(f"Wrote {out_path}")


# ---------------------------------------------------------------------------
# Spec-driven renderer
# ---------------------------------------------------------------------------

# Named color palette accessible from a deck spec JSON (avoids requiring
# every spec to embed hex strings).
_NAMED_COLORS = {
    "NAVY": NAVY,
    "INK": INK,
    "WHITE": WHITE,
    "MCK_BLUE": MCK_BLUE,
    "MCK_BLUE_BR": MCK_BLUE_BR,
    "MCK_BLUE_PL": MCK_BLUE_PL,
    "ACCENT_CORAL": ACCENT_CORAL,
    "GREY_FOOT": GREY_FOOT,
    "GREY_RULE": GREY_RULE,
    "GREY_LIGHT": GREY_LIGHT,
    "RED":   RGBColor(0xC0, 0x39, 0x2B),
    "GREEN": RGBColor(0x1F, 0x7A, 0x44),
    "AMBER": RGBColor(0xE6, 0x9A, 0x29),
}


def _resolve_color(c):
    """Accept a named color, hex string '#RRGGBB', or RGBColor instance."""
    if isinstance(c, RGBColor):
        return c
    if isinstance(c, str):
        if c.startswith("#") and len(c) == 7:
            return RGBColor(int(c[1:3], 16), int(c[3:5], 16), int(c[5:7], 16))
        if c in _NAMED_COLORS:
            return _NAMED_COLORS[c]
    return INK   # safe fallback


def _resolve_table_rows(rows):
    """Convert spec-style rows (lists of strings or [text, color, kind])
    into the tuple shape add_table_slide expects."""
    out = []
    for row in rows:
        new_row = []
        for cell in row:
            if isinstance(cell, list):
                if len(cell) == 3:
                    text, color, kind = cell
                    new_row.append((text, _resolve_color(color), kind))
                elif len(cell) == 2:
                    text, color = cell
                    new_row.append((text, _resolve_color(color)))
                else:
                    new_row.append(cell[0])
            else:
                new_row.append(cell)
        out.append(new_row)
    return out


def _inches(v):
    """Convert a spec width value (number = inches, or string with unit)."""
    if isinstance(v, (int, float)):
        return Inches(v)
    return v


def build_from_spec(spec, out_path):
    """Render a deck from a JSON spec dict. See deck-spec.md for schema."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    meta = spec.get("meta", {})
    brand = meta.get("brand", "Autodecision")

    for slide in spec["slides"]:
        t = slide["type"]

        if t == "title_cover":
            add_title_cover(
                prs,
                title=slide.get("title", meta.get("title", "")),
                subtitle=slide.get("subtitle", meta.get("subtitle", "")),
                footer_org=slide.get("brand", brand),
                date=slide.get("date", meta.get("date", "")),
            )

        elif t == "toc_dark":
            items = [(it["label"], it["pages"])
                     for it in slide.get("items", spec.get("toc", []))]
            add_toc_dark(prs, items, page_num=slide["page_num"])

        elif t == "section_divider":
            add_section_divider(
                prs,
                section_num=slide["section_num"],
                section_title=slide["section_title"],
                lead_question=slide["lead_question"],
                page_num=slide["page_num"],
            )

        elif t == "action_text":
            add_action_text_slide(
                prs,
                prefix=slide.get("prefix"),
                title=slide["title"],
                bullets=slide["bullets"],
                source=slide.get("source", ""),
                page_num=slide["page_num"],
                subtitle=slide.get("subtitle"),
            )

        elif t == "two_column":
            add_two_column_slide(
                prs,
                prefix=slide.get("prefix"),
                title=slide["title"],
                left_header=slide["left_header"],
                left_items=slide["left_items"],
                right_header=slide["right_header"],
                right_items=slide["right_items"],
                source=slide.get("source", ""),
                page_num=slide["page_num"],
            )

        elif t == "three_column":
            cols = [(c["header"], _resolve_color(c["color"]), c["items"])
                    for c in slide["columns"]]
            add_three_column_grid(
                prs,
                prefix=slide.get("prefix"),
                title=slide["title"],
                columns=cols,
                source=slide.get("source", ""),
                page_num=slide["page_num"],
            )

        elif t == "table":
            col_widths = [_inches(w) for w in slide["col_widths"]]
            rows = _resolve_table_rows(slide["rows"])
            add_table_slide(
                prs,
                prefix=slide.get("prefix"),
                title=slide["title"],
                subtitle=slide.get("subtitle"),
                headers=slide["headers"],
                rows=rows,
                col_widths=col_widths,
                source=slide.get("source", ""),
                page_num=slide["page_num"],
            )

        elif t == "recommendation":
            fields = [(f["label"], f["value"], _resolve_color(f["color"]))
                      for f in slide["fields"]]
            add_recommendation_slide(
                prs,
                prefix=slide.get("prefix"),
                title=slide["title"],
                action=slide["action"],
                fields=fields,
                source=slide.get("source", ""),
                page_num=slide["page_num"],
            )

        elif t == "matrix_2x2":
            items = []
            for it in slide["items"]:
                # color names → hex; spec uses string colors throughout
                color = it.get("color", "#054AD8")
                if color in _NAMED_COLORS:
                    rgb = _NAMED_COLORS[color]
                    color = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                items.append({**it, "color": color})
            add_matrix_2x2_slide(
                prs,
                prefix=slide.get("prefix"),
                title=slide["title"],
                axes=slide["axes"],
                items=items,
                quadrant_labels=slide.get("quadrant_labels"),
                source=slide.get("source", ""),
                page_num=slide["page_num"],
            )

        elif t == "radar":
            personas = []
            for p in slide["personas"]:
                color = p["color"]
                if color in _NAMED_COLORS:
                    rgb = _NAMED_COLORS[color]
                    color = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                personas.append({**p, "color": color})
            add_radar_slide(
                prs,
                prefix=slide.get("prefix"),
                title=slide["title"],
                axes=slide["axes"],
                personas=personas,
                commentary_lines=slide.get("commentary_lines"),
                source=slide.get("source", ""),
                page_num=slide["page_num"],
            )

        elif t == "chart_bar":
            chart_buf = make_effects_bar_chart(
                slide["effects"],
                w=slide.get("chart_w_inch", 8.0),
                h=slide.get("chart_h_inch", 4.6),
            )
            add_chart_slide(
                prs,
                prefix=slide.get("prefix"),
                title=slide["title"],
                chart_buf=chart_buf,
                commentary_lines=slide.get("commentary_lines"),
                source=slide.get("source", ""),
                page_num=slide["page_num"],
                chart_w=Inches(slide.get("chart_w_inch", 8.0)),
            )

        elif t == "chart_tornado":
            chart_buf = make_sensitivity_tornado(
                slide["assumptions"],
                w=slide.get("chart_w_inch", 8.0),
                h=slide.get("chart_h_inch", 4.6),
            )
            add_chart_slide(
                prs,
                prefix=slide.get("prefix"),
                title=slide["title"],
                chart_buf=chart_buf,
                commentary_lines=slide.get("commentary_lines"),
                source=slide.get("source", ""),
                page_num=slide["page_num"],
                chart_w=Inches(slide.get("chart_w_inch", 8.0)),
            )

        else:
            raise ValueError(f"Unknown slide type: {t!r}")

    prs.save(out_path)
    print(f"Wrote {out_path}")


def main():
    import argparse, json
    ap = argparse.ArgumentParser(description="Render a Decision Brief deck.")
    ap.add_argument("--spec", help="Path to deck-spec.json")
    ap.add_argument("--out", help="Path to output PPTX (default /tmp/adobe-deck.pptx)")
    ap.add_argument("positional", nargs="?", help="Deprecated; pass --out instead")
    args = ap.parse_args()

    out = Path(args.out or args.positional or "/tmp/adobe-deck.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)

    if args.spec:
        with open(args.spec) as f:
            spec = json.load(f)
        build_from_spec(spec, out)
    else:
        # Default: render the hardcoded Adobe deck (regression test path)
        build_adobe_deck(out)


if __name__ == "__main__":
    main()
